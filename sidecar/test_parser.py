"""pptx_parser 자체 검증 — python-pptx로 더미 PPTX를 만들어 추출을 확인한다.

실행: python sidecar/test_parser.py  (사전: pip install python-pptx)
외부 PPTX 파일 없이 파싱 로직(제목·본문·노트·표·빈 슬라이드·순서)을 검증한다.
"""
from __future__ import annotations

import os
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches

from pptx_parser import extract_slides

PASS = 0
FAIL = 0


def check(name: str, condition: bool, detail: str = "") -> None:
    global PASS, FAIL
    if condition:
        PASS += 1
        print("PASS " + name)
    else:
        FAIL += 1
        print("FAIL " + name + (" — " + detail if detail else ""))


def _build_sample(path: str) -> None:
    presentation = Presentation()

    # 슬라이드 1: 제목 + 본문 + 발표자 노트 (레이아웃 1 = 제목+내용)
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "첫 번째 제목"
    slide1.placeholders[1].text = "본문 첫 줄\n본문 둘째 줄"
    slide1.notes_slide.notes_text_frame.text = "이 슬라이드의 발표자 노트"

    # 슬라이드 2: 제목 + 표 (노트 없음)
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])  # 제목만
    slide2.shapes.title.text = "표가 있는 슬라이드"
    rows, cols = 2, 2
    table_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(4), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "값"
    table.cell(1, 0).text = "노아"
    table.cell(1, 1).text = "동반자"

    # 슬라이드 3: 빈 슬라이드(레이아웃 6 = 완전 빈) — 텍스트 없음
    presentation.slides.add_slide(presentation.slide_layouts[6])

    presentation.save(path)


def main() -> int:
    temp_dir = tempfile.mkdtemp(prefix="aivisor-test-")
    pptx_path = os.path.join(temp_dir, "sample.pptx")
    _build_sample(pptx_path)

    slides = extract_slides(pptx_path)

    check("슬라이드 수 = 3 (빈 슬라이드 포함, 거르기는 renderer 책임)", len(slides) == 3, "got %d" % len(slides))

    s1 = slides[0]
    check("S1 번호 = 1", s1["number"] == 1)
    check("S1 제목 추출", s1["title"] == "첫 번째 제목", repr(s1["title"]))
    check("S1 본문에 두 줄 포함", "본문 첫 줄" in s1["bodyText"] and "본문 둘째 줄" in s1["bodyText"], repr(s1["bodyText"]))
    check("S1 제목이 본문에 중복되지 않음", "첫 번째 제목" not in s1["bodyText"], repr(s1["bodyText"]))
    check("S1 발표자 노트 추출", s1["speakerNotes"] == "이 슬라이드의 발표자 노트", repr(s1["speakerNotes"]))

    s2 = slides[1]
    check("S2 제목 추출", s2["title"] == "표가 있는 슬라이드")
    check("S2 표 셀 텍스트 추출", "노아" in s2["bodyText"] and "동반자" in s2["bodyText"], repr(s2["bodyText"]))
    check("S2 노트 없음 → 빈 문자열", s2["speakerNotes"] == "")

    s3 = slides[2]
    check("S3 빈 슬라이드 → 제목·본문 빈 문자열", s3["title"] == "" and s3["bodyText"] == "", repr(s3))

    print("\n결과: %d PASS / %d FAIL (총 %d)" % (PASS, FAIL, PASS + FAIL))
    return 0 if FAIL == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
