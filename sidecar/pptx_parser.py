"""PPTX 슬라이드별 텍스트·발표자 노트 추출 (python-pptx).

이미지화는 이 모듈의 책임이 아니다 — python-pptx로는 슬라이드 이미지화가 안 되는
것이 알려진 한계다. 슬라이드 렌더링은 slide_renderer.py(LibreOffice)가 담당한다.

반환 dict는 renderer/src/presentation/slides.ts의 Slide와 호환된다:
  {"number": int, "title": str, "bodyText": str, "speakerNotes": str}
"""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_leaf_shapes(shapes):
    """그룹 도형을 재귀적으로 펼쳐 말단 도형만 내보낸다."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _table_text(table) -> str:
    """표 도형의 셀 텍스트를 행 단위로 모은다(빈 셀·빈 행은 건너뛴다)."""
    lines = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        joined = " | ".join(cell for cell in cells if cell)
        if joined:
            lines.append(joined)
    return "\n".join(lines)


def _extract_one(slide, number: int) -> dict:
    title_shape = slide.shapes.title
    title = (title_shape.text or "").strip() if title_shape is not None else ""
    # python-pptx는 같은 도형이라도 접근마다 새 래퍼 객체를 줄 수 있어 `is` 동일성
    # 비교가 빗나간다 → 안정적인 shape_id로 제목 도형을 식별한다(본문 중복 방지).
    title_shape_id = title_shape.shape_id if title_shape is not None else None

    body_parts = []
    for shape in _iter_leaf_shapes(slide.shapes):
        if title_shape_id is not None and shape.shape_id == title_shape_id:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                body_parts.append(text)
        elif getattr(shape, "has_table", False) and shape.has_table:
            text = _table_text(shape.table)
            if text:
                body_parts.append(text)

    notes = ""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is not None:
            notes = (notes_frame.text or "").strip()

    return {
        "number": number,
        "title": title,
        "bodyText": "\n".join(body_parts),
        "speakerNotes": notes,
    }


def extract_slides(pptx_path: str) -> list:
    """PPTX 파일에서 슬라이드별 텍스트·노트를 추출한다. 슬라이드 순서를 보존한다.

    **모든** 슬라이드를 반환한다 — 제목·본문이 비어 있어도 거르지 않는다. "발표
    가능한 슬라이드만" 거르는 것은 렌더러(normalizeSlideDeck)의 책임이다(파서는
    "파일에 무엇이 있는가", 렌더러는 "무엇을 발표할 수 있는가"). 여기서 빈 슬라이드를
    거르면 양쪽이 이중으로 거르거나 번호가 어긋난다.
    """
    presentation = Presentation(pptx_path)
    return [_extract_one(slide, index + 1) for index, slide in enumerate(presentation.slides)]
